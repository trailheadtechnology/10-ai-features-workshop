"""Build the themed .pptx decks from the outlines.

    python3 build.py            # writes pptx/*.pptx from outlines/*.md

Each outline is markdown; `## Title` starts a slide (an optional leading
`N.` is ignored). A `[marker]` before the title picks the slide kind:

    [title]    clone the template's title slide (banner, headshot); the
               outline title goes in the big text box, the first bullet in
               the subtitle box, any second bullet as a smaller third line
    [promise]  clone the template's full-bleed blue promise slide; the
               outline title is the statement (**bold** spans honored)
    [about]    clone the template's about-me slide (bio, QR, free offer)
    [thanks]   clone the template's closing slide ("Thanks! Questions?")
    [section]  Section Header layout, brand-blue title, logo bottom-right
    [deck-title]  a module deck's opener: banner, headshot, and logo around
               "Building Pragmatic AI", with the outline title (the module
               name) in the small bottom placeholder — J's M1 title design
    [demo]     full-bleed blue slide reading DEMO (or the outline title)
    [define]   a feature-definition slide: the outline title is the one-line
               definition (40pt centered, **spans** bold in brand blue), an
               `Icon (name):` line puts that icon centered above it
    [big]      one big centered statement on a blank slide
    [static]   ordinary content slide, everything shown at once

An unmarked slide is an ordinary content slide, and its list is revealed
item by item across consecutive slides (advance-to-animate): `- ` lines are
bullets (two-space indent nests one level), an `Icon (name): Label` line adds
one cell to a centered icon grid (assets/icons/<name>.png above a bold
centered label, logo bottom-right, one cell per advance — the Power of Ten
"Patterns in Failures" style), a `Flow: A -> B -> C` line draws
a left-to-right box-and-arrow diagram revealed box by box (`Flow (label):`
puts a label to its left; several Flow lines stack as rows and reveal row by
row), and any other paragraph is body text. `**bold**` and `` `code` ``
render as bold and Consolas runs. A `Notes:` paragraph becomes the speaker
notes; everything after it up to the next `##` is also notes (that is where
the per-feature demo scripts live). In a reveal sequence a line that is exactly
`~` splits the notes: the chunk before the first `~` goes on the first reveal
slide, the next chunk on the second, and so on, so each slide carries only the
new thing to say. Without any `~` every slide of the sequence carries the same
notes. Everything before the first `##` is the outline's
own preamble (module runsheet) and is skipped.

The theme comes from template.pptx (extracted from the Power of Ten deck:
masters, layouts, and the four boilerplate slides) plus assets/logo.png.
Requires python-pptx.
"""

import copy
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

HERE = Path(__file__).parent
OUTLINES = HERE / "outlines"
OUT = HERE / "pptx"
TEMPLATE = HERE / "template.pptx"
LOGO = HERE / "assets" / "logo.png"
ICON_DIR = HERE / "assets" / "icons"

BRAND = RGBColor(0x01, 0x7A, 0xB4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GRAY = RGBColor(0xD9, 0xD9, 0xD9)
DIFFICULTY = {"easy": 1, "medium": 2, "hard": 3}
DIFFICULTY_COLOR = {"easy": BRAND,
                    "medium": RGBColor(0xFF, 0xC0, 0x00),
                    "hard": RGBColor(0xC0, 0x00, 0x00)}

# template slide order (see make-template provenance in git history)
TPL = {"title": 0, "promise": 1, "about": 2, "thanks": 3}
REPO_OLD, REPO_NEW = "power-of-ten", "10-ai-features-workshop"

DECKS = [
    # M0-opening.pptx and M1-understanding.pptx are hand-tuned by J; M1 is the
    # design master the builder imitates. Diff and port any of J's edits into
    # this script BEFORE re-enabling either line — a rebuild overwrites them.
    # ("M0-opening.md", "M0-opening.pptx"),
    # ("M1-understanding.md", "M1-understanding.pptx"),
    ("M2-finding.md", "M2-finding.pptx"),
    ("M3-deciding.md", "M3-deciding.pptx"),
    ("M4-doing.md", "M4-doing.pptx"),
    ("M5-closing.md", "M5-closing.pptx"),
]


def clean(text):
    """Drop single-star emphasis; keep ** and ` for run styling."""
    text = text.replace("**", "\x00")
    text = re.sub(r"(?<!\w)\*(.+?)\*(?!\w)", r"\1", text)
    return text.replace("\x00", "**").strip()


def plain(text):
    return text.replace("**", "").replace("`", "")


RICH_RE = re.compile(r"(\*\*.+?\*\*|`[^`]+`)")


def add_rich(para, text, size=None, color=None, bold=False, font=None, italic=False):
    for tok in RICH_RE.split(text):
        if not tok:
            continue
        b, f, t = bold, font, tok
        if tok.startswith("**"):
            t, b = tok[2:-2], True
        elif tok.startswith("`"):
            t, f = tok[1:-1], "Consolas"
        r = para.add_run()
        r.text = t
        if size is not None:
            r.font.size = Pt(size)
        if b:
            r.font.bold = True
        if italic:
            r.font.italic = True
        if f:
            r.font.name = f
        if color is not None:
            r.font.color.rgb = color


def parse(path):
    """Return a list of slides: {kind, title, items:[(level,text)], flows, notes}."""
    slides = []
    cur = None
    in_notes = False
    for raw in path.read_text().splitlines():
        line = raw.rstrip()
        m = re.match(r"^##\s+(?:\[([\w-]+)\]\s*)?(?:\d+\.\s+)?(.*)$", line)
        if m:
            cur = {"kind": m.group(1) or "content", "title": clean(m.group(2)),
                   "items": [], "flows": [], "icons": [], "notes": [], "difficulty": None}
            slides.append(cur)
            in_notes = False
            continue
        if cur is None or not line.strip() or line.strip() == "---":
            continue
        if line.startswith("Notes:"):
            cur["notes"].append(clean(line[len("Notes:"):]))
            in_notes = True
            continue
        if in_notes:
            cur["notes"].append(clean(line))
            continue
        m = re.match(r"^Difficulty:\s*(easy|medium|hard)\s*$", line, re.I)
        if m:
            cur["difficulty"] = m.group(1).lower()
            continue
        m = re.match(r"^Icon\s*\((.*?)\):\s*(.*)$", line)
        if m:
            cur["icons"].append((m.group(1).strip(), plain(clean(m.group(2)))))
            continue
        m = re.match(r"^Flow(?:\s*\((.*?)\))?:\s*(.*)$", line)
        if m:
            boxes = [plain(clean(b)) for b in m.group(2).split("->")]
            cur["flows"].append((plain(clean(m.group(1) or "")), boxes))
            continue
        m = re.match(r"^(\s*)-\s+(.*)$", line)
        if m:
            level = 1 if len(m.group(1)) >= 2 else 0
            cur["items"].append((level, clean(m.group(2))))
            continue
        cur["items"].append((0, clean(line)))
    for s in slides:
        s["notes"] = "\n".join(s["notes"]).strip()
    return slides


# --- template-slide plumbing ---------------------------------------------


def retext(tf_or_para, text, keep_run=True):
    """Replace a paragraph's text, preserving the first run's formatting."""
    p = tf_or_para
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p.add_run().text = text
    for br in p._p.findall(qn("a:br")):
        p._p.remove(br)


def drop_extra_paragraphs(tf):
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)


def max_run_size(shape):
    sizes = [r.font.size.pt for p in shape.text_frame.paragraphs
             for r in p.runs if r.font.size]
    return max(sizes) if sizes else 0


def fill_title_slide(slide, s):
    lines = [t for _, t in s["items"]]
    boxes = [sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
    main = max(boxes, key=max_run_size)
    retext(main.text_frame.paragraphs[0], plain(s["title"]))
    drop_extra_paragraphs(main.text_frame)
    subs = sorted((b for b in boxes if b is not main and max_run_size(b) >= 30),
                  key=max_run_size, reverse=True)
    if subs and lines:
        tf = subs[0].text_frame
        retext(tf.paragraphs[0], plain(lines[0]))
        drop_extra_paragraphs(tf)
        for extra in lines[1:]:
            p = tf.add_paragraph()
            add_rich(p, extra, 20)


def fill_promise_slide(slide, s):
    ph = next(sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip())
    tf = ph.text_frame
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for br in p0._p.findall(qn("a:br")):
        p0._p.remove(br)
    drop_extra_paragraphs(tf)
    add_rich(p0, s["title"], 54, color=WHITE)


def swap_repo_url(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if REPO_OLD in r.text:
                    r.text = r.text.replace(REPO_OLD, REPO_NEW)


def set_notes(slide, notes, clear_empty=False):
    """One paragraph per beat; `code` spans (full CLI commands) render in Consolas.

    Empty notes leave the slide alone unless clear_empty is set, which the
    reveal splitter uses so a slide with nothing new to say has blank notes.
    """
    if not notes and not clear_empty:
        return
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    if not notes:
        return
    notes = "\n".join(l for l in notes.split("\n") if l.strip() != "~")
    for i, line in enumerate(notes.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        add_rich(p, line)


def delete_slide(prs, slide):
    lst = prs.slides._sldIdLst
    for sldId in list(lst):
        if prs.part.related_part(sldId.get(qn("r:id"))) is slide.part:
            prs.part.drop_rel(sldId.get(qn("r:id")))
            lst.remove(sldId)
            return


def reorder(prs, ordered):
    lst = prs.slides._sldIdLst
    by_part = {}
    for sldId in list(lst):
        by_part[id(prs.part.related_part(sldId.get(qn("r:id"))))] = sldId
        lst.remove(sldId)
    for slide in ordered:
        lst.append(by_part[id(slide.part)])


# --- generated slides -----------------------------------------------------


def layout(prs, name):
    for l in prs.slide_masters[0].slide_layouts:
        if l.name == name:
            return l
    raise KeyError(name)


def add_logo(slide):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(10.2), Inches(6.55), Inches(2.8), Inches(0.7))


def section_slide(prs, s):
    slide = prs.slides.add_slide(layout(prs, "Section Header"))
    tf = slide.shapes.title.text_frame
    tf.word_wrap = True
    add_rich(tf.paragraphs[0], s["title"], 40, color=BRAND, bold=True)
    # clear the small body placeholder if present
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text_frame.text = ""
    add_logo(slide)
    return slide


def deck_title_slide(prs, s):
    """Module-deck opener matching J's hand-built M1 title slide."""
    slide = prs.slides.add_slide(layout(prs, "Section Header"))
    tf = slide.shapes.title.text_frame
    tf.word_wrap = True
    add_rich(tf.paragraphs[0], "Building Pragmatic AI", 40, bold=True)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text_frame.text = plain(s["title"])
    add_logo(slide)
    slide.shapes.add_picture(str(HERE / "assets" / "banner.png"),
                             0, Inches(0.01), Inches(13.33), Inches(3.48))
    slide.shapes.add_picture(str(HERE / "assets" / "headshot.png"),
                             Inches(10.14), Inches(4.0), Inches(2.37), Inches(2.37))
    return slide


def define_slide(prs, s):
    """One-sentence feature definition: icon on top, 40pt centered sentence,
    bold spans in brand blue, logo bottom-right."""
    slide = prs.slides.add_slide(layout(prs, "Blank"))
    if s["icons"]:
        icon_h = Inches(1.0)
        pic = slide.shapes.add_picture(str(ICON_DIR / f"{s['icons'][0][0]}.png"),
                                       0, Inches(1.5), height=icon_h)
        pic.left = int((prs.slide_width - pic.width) / 2)
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(3.1), Inches(11.3), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for tok in RICH_RE.split(s["title"]):
        if not tok:
            continue
        r = p.add_run()
        if tok.startswith("**"):
            r.text = tok[2:-2]
            r.font.bold = True
            r.font.color.rgb = BRAND
        elif tok.startswith("`"):
            r.text = tok[1:-1]
            r.font.name = "Consolas"
        else:
            r.text = tok
        r.font.size = Pt(40)
    add_logo(slide)
    return slide


def demo_slide(prs, s):
    slide = prs.slides.add_slide(layout(prs, "Blank"))
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = BRAND
    rect.line.fill.background()
    rect.shadow.inherit = False
    tb = slide.shapes.add_textbox(Inches(1), Inches(3.19), Inches(11.3), Inches(1.11))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_rich(p, s["title"] or "**DEMO**", 60, color=WHITE)
    add_logo(slide)
    return slide


def big_slide(prs, s, items_upto=None):
    slide = prs.slides.add_slide(layout(prs, "Blank"))
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.9))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_rich(p, s["title"], 54)
    for _, text in s["items"][: items_upto if items_upto is not None else len(s["items"])]:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        add_rich(p, text, 28)
    return slide


def draw_meter(slide, left, top, level_name):
    """Three ascending bars, filled to the difficulty level, level word beneath."""
    level = DIFFICULTY[level_name]
    bar_w, gap = Inches(0.28), Inches(0.10)
    heights = [Inches(0.28), Inches(0.50), Inches(0.72)]
    bottom = top + heights[-1]
    for i, h in enumerate(heights):
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     left + i * (bar_w + gap), bottom - h, bar_w, h)
        bar.adjustments[0] = 0.25
        bar.fill.solid()
        bar.fill.fore_color.rgb = DIFFICULTY_COLOR[level_name] if i < level else GRAY
        bar.line.fill.background()
        bar.shadow.inherit = False
    total_w = bar_w * 3 + gap * 2
    tb = slide.shapes.add_textbox(left - Inches(0.3), bottom + Inches(0.06),
                                  total_w + Inches(0.6), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_rich(p, level_name.capitalize(), 14, bold=True)


def draw_icons(slide, prs, icons, upto=None):
    """Grid of icon-above-label cells spread across the slide width, one cell
    per reveal step. Geometry matches J's hand-tuned M1 This Module slide:
    cell centers from 2.225" to 11.325", icons at y=3.27", 24pt bold labels."""
    shown = icons[: upto if upto is not None else len(icons)]
    n = len(icons)
    per_row = n if n <= 5 else -(-n // 2)
    cell_w = Inches(2.35)
    icon_h = Inches(0.95)
    row_h = Inches(2.3)
    c_first, c_last = Inches(2.225), Inches(11.325)
    top0 = Inches(3.27) if n <= per_row else Inches(2.0)
    for i, (name, label) in enumerate(shown):
        row, col = divmod(i, per_row)
        in_row = min(per_row, n - row * per_row)
        if in_row == 1:
            center = (c_first + c_last) / 2
        else:
            center = c_first + (c_last - c_first) * col // (in_row - 1)
        y = top0 + row_h * row
        pic = slide.shapes.add_picture(str(ICON_DIR / f"{name}.png"), 0, y, height=icon_h)
        pic.left = int(center - pic.width / 2)
        tb = slide.shapes.add_textbox(int(center - cell_w / 2), y + icon_h + Inches(0.16),
                                      cell_w, Inches(0.85))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        add_rich(p, label, 24, bold=True)
    add_logo(slide)


def draw_flows(slide, prs, flows, top, boxes_upto=None):
    """Rows of rounded boxes joined by arrows; brand-blue outline, black text.
    boxes_upto limits how many boxes of a single-row flow are drawn.
    Returns the bottom edge so body text can start below the diagram."""
    margin = Inches(0.67)
    label_w = Inches(1.6) if any(lbl for lbl, _ in flows) else 0
    left0 = margin + label_w
    avail = prs.slide_width - margin - left0
    gap_h = Inches(0.3)
    y = top
    for label, boxes in flows:
        n = len(boxes)
        shown = n if boxes_upto is None else min(n, boxes_upto)
        arrow_w = Inches(0.45)
        box_w = (avail - arrow_w * (n - 1)) / n
        pt = 15 if n <= 4 else 13
        cpl = max(6, int((box_w / 914400 - 0.16) * (8 if pt == 15 else 9)))
        lines = max(-(-len(b) // cpl) for b in boxes)
        row_h = max(Inches(0.9), Inches(0.35 + 0.26 * lines))
        if label:
            tb = slide.shapes.add_textbox(margin, y, label_w - Inches(0.1), row_h)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.text = label
            for r in tf.paragraphs[0].runs:
                r.font.size = Pt(14)
                r.font.italic = True
        x = left0
        for i, text in enumerate(boxes[:shown]):
            shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, row_h)
            shp.fill.background()
            shp.line.color.rgb = BRAND
            shp.line.width = Pt(1.75)
            shp.shadow.inherit = False
            tf = shp.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = Inches(0.08)
            tf.text = text
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            for r in para.runs:
                r.font.size = Pt(pt)
                r.font.color.rgb = BLACK
            x += box_w
            if i < shown - 1:
                ar = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, x + Inches(0.06),
                    y + row_h / 2 - Inches(0.16), arrow_w - Inches(0.12), Inches(0.32))
                ar.fill.solid()
                ar.fill.fore_color.rgb = BRAND
                ar.line.fill.background()
                ar.shadow.inherit = False
                x += arrow_w
        y += row_h + gap_h
    return y - gap_h


def content_slide(prs, s, flow_upto=None, items_upto=None, show_difficulty=False):
    """One rendering of a content slide: flows (possibly partial), then items."""
    slide = prs.slides.add_slide(layout(prs, "Title and Content"))
    slide.shapes.title.text_frame.word_wrap = True
    retext(slide.shapes.title.text_frame.paragraphs[0], plain(s["title"]))
    body = next(ph for ph in slide.placeholders if ph.placeholder_format.idx == 1)
    if s.get("icons"):
        body._element.getparent().remove(body._element)
        draw_icons(slide, prs, s["icons"], upto=flow_upto)
        return slide
    tf = body.text_frame
    tf.word_wrap = True
    n_total = len(s["items"])
    base = 28 if n_total <= 6 else 24
    body_top = Inches(2.1)
    if s["flows"]:
        flows = s["flows"]
        if flow_upto is not None:
            if len(flows) > 1:
                flows = flows[:flow_upto]
                bottom = draw_flows(slide, prs, flows, Inches(1.9))
            else:
                bottom = draw_flows(slide, prs, flows, Inches(1.9), boxes_upto=flow_upto)
        else:
            bottom = draw_flows(slide, prs, flows, Inches(1.9))
        body_top = bottom + Inches(0.25)
    margin = Inches(0.67)
    body.left, body.top = margin, body_top
    body.width, body.height = prs.slide_width - 2 * margin, Inches(7.15) - body_top
    shown = s["items"][: items_upto if items_upto is not None else n_total]
    title_plain = plain(s["title"])
    lab_slide = title_plain.startswith("Lab ")
    first = True
    if lab_slide and title_plain.endswith("Debrief"):
        tf.paragraphs[0].add_run().text = ""  # J's spacer line under the title
        first = False
    for level, text in shown:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        size = base if level == 0 else base - 3
        if lab_slide:  # J's lab-slide hierarchy: big pick lines, medium done-check
            size = 36 if text.startswith(("⭐", "⛰")) else 32 if text.startswith("✅") else size
            if text.startswith("✅"):  # two spacer lines above the done-check
                # the paragraph created above becomes the first spacer
                tf.add_paragraph()
                p = tf.add_paragraph()
        add_rich(p, text, size)
    if show_difficulty and s.get("difficulty"):
        tf.add_paragraph()  # J's spacer line so the label aligns with the meter
        p = tf.add_paragraph()
        add_rich(p, "**Difficulty:** ", base)
        draw_meter(slide, Inches(3.07), Inches(4.82), s["difficulty"])
    if s.get("difficulty"):
        for p in tf.paragraphs:
            p.space_after = Pt(14)
    return slide


def emit(prs, s):
    """Render one outline slide into one or more pptx slides; return them."""
    out = []
    kind = s["kind"]
    if kind == "section":
        out.append(section_slide(prs, s))
    elif kind == "deck-title":
        out.append(deck_title_slide(prs, s))
    elif kind == "demo":
        out.append(demo_slide(prs, s))
    elif kind == "define":
        out.append(define_slide(prs, s))
    elif kind == "big":
        out.append(big_slide(prs, s))
    elif kind == "static":
        out.append(content_slide(prs, s, show_difficulty=True))
    elif s.get("icons"):
        for k in range(1, len(s["icons"]) + 1):
            out.append(content_slide(prs, s, flow_upto=k))
    else:
        flows, items = s["flows"], s["items"]
        flow_steps = (len(flows) if len(flows) > 1 else (len(flows[0][1]) if flows else 0))
        reveal = flow_steps + len(items) > 1
        if not reveal:
            out.append(content_slide(prs, s))
        else:
            for k in range(1, flow_steps + 1):
                out.append(content_slide(prs, s, flow_upto=k, items_upto=0))
            for k in range(0 if flow_steps else 1, len(items) + 1):
                if k == 0:
                    continue
                out.append(content_slide(prs, s, items_upto=k))
            if s.get("difficulty"):
                out.append(content_slide(prs, s, show_difficulty=True))
    for slide, notes in zip(out, split_notes(s["notes"], len(out))):
        set_notes(slide, notes, clear_empty=len(out) > 1)
    return out


def split_notes(notes, n):
    """Distribute a Notes block across the n slides of a reveal sequence.

    A line that is exactly `~` separates chunks; chunk k goes to reveal slide k,
    so each slide's notes hold only the new thing to say. With no `~` the whole
    block goes on every slide (the old behavior). Fewer chunks than slides leaves
    the trailing slides without notes; extra chunks fold into the last slide.
    """
    chunks = [c.strip() for c in re.split(r"(?m)^~\s*$", notes or "")]
    if len(chunks) == 1:
        return [notes] * n
    if len(chunks) > n:
        chunks = chunks[:n - 1] + ["\n".join(chunks[n - 1:])]
    return chunks + [""] * (n - len(chunks))


def build(outline, target):
    prs = Presentation(TEMPLATE)
    tpl_slides = list(prs.slides)  # title, promise, about, thanks
    used_tpl = set()
    ordered = []

    for s in parse(outline):
        kind = s["kind"]
        if kind in TPL:
            slide = tpl_slides[TPL[kind]]
            used_tpl.add(TPL[kind])
            if kind == "title":
                fill_title_slide(slide, s)
            elif kind == "promise":
                fill_promise_slide(slide, s)
            else:
                swap_repo_url(slide)
            set_notes(slide, s["notes"])
            ordered.append(slide)
        else:
            ordered.extend(emit(prs, s))

    for i, slide in enumerate(tpl_slides):
        if i not in used_tpl:
            delete_slide(prs, slide)
    reorder(prs, ordered)

    OUT.mkdir(exist_ok=True)
    prs.save(OUT / target)
    return len(prs.slides)


if __name__ == "__main__":
    for src, dst in DECKS:
        n = build(OUTLINES / src, dst)
        print(f"{dst}: {n} slides")
    sys.exit(0)
